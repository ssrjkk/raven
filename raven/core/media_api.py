from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from fastapi import APIRouter, HTTPException, UploadFile
from loguru import logger


def create_media_router(workspace_dir: str | Path = "") -> APIRouter:
    router = APIRouter(prefix="/api/media", tags=["media"])
    ws = Path(workspace_dir) if workspace_dir else Path.cwd() / "workspace"

    @router.post("/generate")
    async def generate_image(prompt: str, size: str = "1024x1024", quality: str = "standard"):
        api_key = os.environ.get("OPENAI_API_KEY", "")
        if not api_key:
            raise HTTPException(400, "OPENAI_API_KEY not configured")
        try:
            import httpx
            async with httpx.AsyncClient(timeout=60) as client:
                resp = await client.post(
                    "https://api.openai.com/v1/images/generations",
                    headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                    json={"model": "dall-e-3", "prompt": prompt, "n": 1, "size": size, "quality": quality},
                )
                resp.raise_for_status()
                data = resp.json()
                urls = [item["url"] for item in data.get("data", []) if item.get("url")]
                if not urls:
                    raise HTTPException(500, "No image URLs returned")
                return {"url": urls[0], "prompt": prompt, "size": size}
        except HTTPException:
            raise
        except Exception as e:
            logger.error("Image generation API error: {}", e)
            raise HTTPException(500, str(e)) from e

    def _confine(p: Path, base: Path) -> Path:
        r = p.expanduser().resolve()
        b = base.resolve()
        try:
            r.relative_to(b)
        except ValueError:
            raise HTTPException(403, f"Path outside workspace: {p}") from None
        return r

    @router.post("/process")
    async def process_image(
        filepath: str,
        resize: str = "",
        crop: str = "",
        rotate: int = 0,
        flip: str = "",
        output_format: str = "",
        quality: int = 85,
    ):
        target = _confine(Path(filepath), ws)
        if not target.exists():
            raise HTTPException(404, f"File not found: {target}")
        try:
            import base64
            from io import BytesIO

            from PIL import Image as PILImage
            src: Any = PILImage.open(target)
            if crop:
                vals = [x.strip() for x in crop.split(",")]
                if len(vals) == 4 and all(v.lstrip("-").isdigit() for v in vals):
                    src = src.crop((int(vals[0]), int(vals[1]), int(vals[2]), int(vals[3])))
            if resize:
                vals = resize.lower().split("x")
                if len(vals) == 2 and vals[0].strip().isdigit() and vals[1].strip().isdigit():
                    src = src.resize((int(vals[0]), int(vals[1])), PILImage.Resampling.LANCZOS)
            if rotate:
                src = src.rotate(rotate, expand=True)
            if flip == "horizontal":
                src = src.transpose(PILImage.FLIP_LEFT_RIGHT)  # type: ignore[attr-defined]
            elif flip == "vertical":
                src = src.transpose(PILImage.FLIP_TOP_BOTTOM)  # type: ignore[attr-defined]
            fmt = output_format.upper() if output_format else (src.format or "PNG")
            if fmt == "JPEG" and src.mode in ("RGBA", "P"):
                src = src.convert("RGB")
            buf = BytesIO()
            src.save(buf, format=fmt, quality=quality)
            b64 = base64.b64encode(buf.getvalue()).decode()
            return {
                "format": fmt.lower(),
                "width": src.size[0],
                "height": src.size[1],
                "bytes": buf.tell(),
                "data_url": f"data:image/{fmt.lower()};base64,{b64}",
            }
        except Exception as e:
            logger.error("Image process error: {}", e)
            raise HTTPException(500, str(e)) from e

    @router.post("/parse")
    async def parse_document(filepath: str, pages: str = ""):
        target = _confine(Path(filepath), ws)
        if not target.exists():
            raise HTTPException(404, f"File not found: {target}")
        ext = target.suffix.lower()
        try:
            text = ""
            meta: dict[str, Any] = {"filename": target.name, "format": ext.lstrip(".")}
            if ext == ".pdf":
                import fitz
                doc = fitz.open(str(target))
                meta["pages"] = len(doc)
                selected: list[int] | range = range(len(doc))
                if pages:
                    pr: list[int] = []
                    for part in pages.split(","):
                        part = part.strip()
                        if "-" in part:
                            s, e = part.split("-", 1)
                            pr.extend(range(int(s.strip()) - 1, int(e.strip())))
                        else:
                            pr.append(int(part) - 1)
                    selected = pr
                chunks: list[str] = []
                for i in selected:
                    if i < len(doc):
                        p = doc[i]
                        t = p.get_text().strip()
                        if t:
                            chunks.append(f"--- Page {i+1} ---\n{t}")
                doc.close()
                text = "\n".join(chunks)
                meta["extracted_pages"] = len(selected)
            elif ext == ".docx":
                from docx import Document
                doc = Document(str(target))
                paras = [p.text for p in doc.paragraphs if p.text.strip()]
                text = "\n".join(paras)
                meta["paragraphs"] = len(paras)
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(target))
                meta["slides"] = len(prs.slides)
                slide_chunks: list[str] = []
                for i, slide in enumerate(prs.slides, 1):
                    lines: list[str] = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    lines.append(para.text)
                    if lines:
                        slide_chunks.append(f"--- Slide {i} ---\n" + "\n".join(lines))
                text = "\n".join(slide_chunks)
            elif ext in (".xlsx", ".xls"):
                from openpyxl import load_workbook
                wb = load_workbook(str(target), read_only=True, data_only=True)
                meta["sheets"] = wb.sheetnames
                sheet_chunks: list[str] = []
                for name in wb.sheetnames:
                    sheet = wb[name]
                    rows: list[str] = []
                    for row in sheet.iter_row():
                        vals = [str(c.value or "") for c in row]
                        rows.append(" | ".join(vals))
                    if rows:
                        sheet_chunks.append(f"--- Sheet: {name} ---\n" + "\n".join(rows))
                wb.close()
                text = "\n".join(sheet_chunks)
            else:
                raise HTTPException(400, f"Unsupported format: {ext}")
            meta["chars"] = len(text)
            return {"text": text[:50000], "truncated": len(text) > 50000, "metadata": meta}
        except HTTPException:
            raise
        except ImportError as e:
            raise HTTPException(500, f"Missing dependency: {e}") from e
        except Exception as e:
            logger.error("Document parse error: {}", e)
            raise HTTPException(500, str(e)) from e

    @router.post("/upload")
    async def upload_file(file: UploadFile):
        if not file.filename:
            raise HTTPException(400, "No filename")
        safe_name = Path(file.filename).name
        if ".." in safe_name or "/" in safe_name or "\\" in safe_name:
            raise HTTPException(400, "Invalid filename")
        ws.mkdir(parents=True, exist_ok=True)
        dest = ws / safe_name
        try:
            content = await file.read()
            dest.write_bytes(content)
            return {"path": str(dest), "size": len(content), "filename": safe_name}
        except Exception as e:
            logger.error("Upload error: {}", e)
            raise HTTPException(500, str(e)) from e

    @router.post("/analyze")
    async def analyze_image(filepath: str, prompt: str = "Describe this image in detail"):
        from raven.tools.media import image_analyze
        result = await image_analyze(filepath, prompt)
        return {"result": result}

    @router.post("/video-info")
    async def video_info_endpoint(filepath: str):
        from raven.tools.media import video_info
        result = await video_info(filepath)
        return {"result": result}

    @router.post("/video-thumbnail")
    async def video_thumbnail_endpoint(filepath: str, time_sec: float = 1.0, size: str = "320x240"):
        from raven.tools.media import video_thumbnail
        result = await video_thumbnail(filepath, time_sec, size)
        return {"result": result}

    @router.post("/video-transcribe")
    async def video_transcribe_endpoint(filepath: str, language: str = ""):
        from raven.tools.media import video_transcribe
        result = await video_transcribe(filepath, "whisper-1", language)
        return {"result": result}

    @router.post("/video-extract-frames")
    async def video_extract_frames_endpoint(filepath: str, interval_sec: float = 5.0, max_frames: int = 10, size: str = "640x480"):
        from raven.tools.media import video_extract_frames
        result = await video_extract_frames(filepath, interval_sec, max_frames, size)
        return {"result": result}

    return router
