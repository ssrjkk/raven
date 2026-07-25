from __future__ import annotations

import asyncio
import base64
import os
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Any

from loguru import logger

from raven.core.task_engine.tool_registry import ToolRegistry, ToolSpec
from raven.tools.file import _confine

_OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
_REPLICATE_API_TOKEN = os.environ.get("REPLICATE_API_TOKEN", "")


async def image_generate(
    prompt: str, size: str = "1024x1024", quality: str = "standard", n: int = 1, model: str = "dall-e-3"
) -> str:
    if model.startswith("dall-e"):
        return await _image_generate_dalle(prompt, size, quality, n)
    if model in ("stable-diffusion", "sdxl", "sd"):
        return await _image_generate_sd(prompt, size)
    return f"[error] Unknown model '{model}'. Supported: dall-e-3, stable-diffusion"


async def _image_generate_dalle(prompt: str, size: str, quality: str, n: int) -> str:
    api_key = _OPENAI_API_KEY
    if not api_key:
        return "[error] OPENAI_API_KEY env var required for DALL-E"
    try:
        import httpx

        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                "https://api.openai.com/v1/images/generations",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": "dall-e-3", "prompt": prompt, "n": n, "size": size, "quality": quality},
            )
            resp.raise_for_status()
            data = resp.json()
            urls = [item["url"] for item in data.get("data", []) if item.get("url")]
            if not urls:
                return "[error] No image URLs returned"
            lines = [f"Generated image(s) via DALL-E for: {prompt}\n"]
            for i, url in enumerate(urls, 1):
                lines.append(f"{i}. ![Generated Image {i}]({url})")
            return "\n".join(lines)
    except Exception as e:
        logger.error("DALL-E generation failed: {}", e)
        return f"[error] DALL-E generation failed: {e}"


async def _image_generate_sd(prompt: str, size: str) -> str:
    api_token = _REPLICATE_API_TOKEN
    if not api_token:
        return "[error] REPLICATE_API_TOKEN env var required for Stable Diffusion (via Replicate)"
    width, height = 1024, 1024
    if size:
        parts = size.lower().split("x")
        if len(parts) == 2:
            import contextlib

            with contextlib.suppress(ValueError):
                width, height = int(parts[0]), int(parts[1])
    try:
        import httpx

        model_version = "stability-ai/stable-diffusion-3.5-large-turbo"
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(
                "https://api.replicate.com/v1/predictions",
                headers={"Authorization": f"Bearer {api_token}", "Content-Type": "application/json"},
                json={
                    "version": model_version,
                    "input": {
                        "prompt": prompt,
                        "width": width,
                        "height": height,
                        "num_outputs": 1,
                    },
                },
            )
            resp.raise_for_status()
            prediction = resp.json()
            get_url = prediction["urls"]["get"]
            for _ in range(30):
                await asyncio.sleep(2)
                poll = await client.get(get_url, headers={"Authorization": f"Bearer {api_token}"})
                poll.raise_for_status()
                status = poll.json()
                if status["status"] == "succeeded":
                    output = status.get("output", [])
                    if isinstance(output, list) and output:
                        lines = [f"Generated image(s) via Stable Diffusion for: {prompt}\n"]
                        for i, url in enumerate(output, 1):
                            lines.append(f"{i}. ![Generated Image {i}]({url})")
                        return "\n".join(lines)
                    return "[error] No output from Stable Diffusion"
                if status["status"] == "failed":
                    return f"[error] Stable Diffusion prediction failed: {status.get('error', 'unknown')}"
            return "[error] Stable Diffusion prediction timed out"
    except Exception as e:
        logger.error("Stable Diffusion generation failed: {}", e)
        return f"[error] Stable Diffusion generation failed: {e}"


async def image_edit(
    filepath: str,
    resize: str = "",
    crop: str = "",
    rotate: int = 0,
    flip: str = "",
    format: str = "",
    quality: int = 85,
    output: str = "",
) -> str:
    try:
        from PIL import Image
    except ImportError as e:
        return f"[error] Pillow not available: {e}"
    try:
        p = _confine(filepath)

        def _process() -> str:
            src: Any = Image.open(p)
            orig_size = src.size
            orig_mode = src.mode
            operations: list[str] = []

            if crop:
                vals = [x.strip() for x in crop.split(",")]
                if len(vals) == 4:
                    left, upper, right, lower = int(vals[0]), int(vals[1]), int(vals[2]), int(vals[3])
                    if right > left and lower > upper:
                        src = src.crop((left, upper, right, lower))
                        operations.append(f"crop({crop})")
                    else:
                        return f"[error] Invalid crop rectangle: {crop} (right must be > left, lower must be > upper)"
                else:
                    return f"[error] crop expects 4 comma-separated values (left,upper,right,lower), got {len(vals)}"

            if resize:
                vals = resize.lower().split("x")
                if len(vals) == 2 and vals[0].strip().isdigit() and vals[1].strip().isdigit():
                    w, h = int(vals[0]), int(vals[1])
                    src = src.resize((w, h), Image.Resampling.LANCZOS)
                    operations.append(f"resize({resize})")
                else:
                    return f"[error] resize expects WxH format (e.g. 800x600), got '{resize}'"

            if rotate:
                src = src.rotate(rotate, expand=True)
                operations.append(f"rotate({rotate})")

            if flip:
                if flip == "horizontal":
                    src = src.transpose(Image.FLIP_LEFT_RIGHT)  # type: ignore[attr-defined]
                    operations.append("flip(horizontal)")
                elif flip == "vertical":
                    src = src.transpose(Image.FLIP_TOP_BOTTOM)  # type: ignore[attr-defined]
                    operations.append("flip(vertical)")
                else:
                    return f"[error] flip expects 'horizontal' or 'vertical', got '{flip}'"

            if format:
                fmt = format.upper()
                if fmt not in ("JPEG", "PNG", "GIF", "WEBP", "BMP", "TIFF"):
                    return f"[error] Unsupported output format: {format}"
                if fmt == "JPEG" and src.mode in ("RGBA", "P"):
                    src = src.convert("RGB")
            else:
                fmt = src.format or "PNG"

            if output:
                out_path = Path(output)
                out_path.parent.mkdir(parents=True, exist_ok=True)
                src.save(out_path, format=fmt, quality=quality)
            else:
                suffix = fmt.lower()
                if suffix == "jpeg":
                    suffix = "jpg"
                out_path = Path(filepath).parent / f"{Path(filepath).stem}_edited.{suffix}"
                src.save(out_path, format=fmt, quality=quality)

            buf = BytesIO()
            src.save(buf, format=fmt, quality=quality)
            b64 = base64.b64encode(buf.getvalue()).decode()

            lines = [
                f"Image edited: {Path(filepath).name} \u2192 {out_path.name}",
                f"- Original: {orig_size[0]}x{orig_size[1]}, {orig_mode}",
                f"- Final: {src.size[0]}x{src.size[1]}, {src.mode}",
            ]
            if operations:
                lines.append(f"- Operations: {', '.join(operations)}")
            lines.append(f"- Size: {buf.tell()} bytes")
            lines.append(f"![{out_path.name}](data:image/{fmt.lower()};base64,{b64})")
            return "\n".join(lines)

        return await asyncio.to_thread(_process)
    except Exception as e:
        logger.error("Image edit failed: {}", e)
        return f"[error] Image processing failed: {e}"


async def document_parse(filepath: str, pages: str = "") -> str:
    path = _confine(filepath)
    exists = await asyncio.to_thread(path.exists)
    if not exists:
        return f"[error] File not found: {filepath}"

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return await asyncio.to_thread(_parse_pdf, path, pages)
        if ext == ".docx":
            return await asyncio.to_thread(_parse_docx, path)
        if ext == ".pptx":
            return await asyncio.to_thread(_parse_pptx, path)
        if ext in (".xlsx", ".xls"):
            return await asyncio.to_thread(_parse_xlsx, path)
        return f"[error] Unsupported format: {ext} (supported: .pdf, .docx, .pptx, .xlsx)"
    except ImportError as e:
        return f"[error] Missing dependency for {ext}: {e}"
    except Exception as e:
        logger.error("Document parse failed for {}: {}", filepath, e)
        return f"[error] Failed to parse {path.name}: {e}"


def _parse_pdf(path: Path, pages: str = "") -> str:
    import fitz

    doc = fitz.open(str(path))
    total_pages = len(doc)
    page_nums: list[int] = []
    if pages:
        for part in pages.split(","):
            part = part.strip()
            if "-" in part:
                start, end = part.split("-", 1)
                page_nums.extend(range(int(start.strip()) - 1, int(end.strip())))
            else:
                page_nums.append(int(part) - 1)
    else:
        page_nums = list(range(total_pages))

    lines = [f"Document: {path.name} ({total_pages} pages)"]
    char_count = 0
    for i in page_nums:
        if i >= total_pages:
            continue
        page = doc[i]
        text = page.get_text().strip()
        if text:
            char_count += len(text)
            lines.append(f"\n--- Page {i + 1} ---\n{text[:2000]}")
            if len(text) > 2000:
                lines.append(f"[... truncated, {len(text)} chars total on this page]")
    doc.close()
    lines.append(f"\n--- Extracted {char_count} chars from {len(page_nums)} pages ---")
    return "\n".join(lines)


def _parse_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    lines = [f"Document: {path.name}"]
    char_count = 0
    for para in doc.paragraphs:
        if para.text.strip():
            char_count += len(para.text)
            lines.append(para.text)
    lines.append(f"\n--- Extracted {char_count} chars from {len(doc.paragraphs)} paragraphs ---")
    return "\n".join(lines)


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines = [f"Presentation: {path.name} ({len(prs.slides)} slides)"]
    char_count = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_lines.append(para.text)
        if slide_lines:
            text = "\n".join(slide_lines)
            char_count += len(text)
            lines.append(f"\n--- Slide {i} ---\n{text}")
    lines.append(f"\n--- Extracted {char_count} chars from {len(prs.slides)} slides ---")
    return "\n".join(lines)


def _parse_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = [f"Workbook: {path.name} ({len(wb.sheetnames)} sheets)"]
    char_count = 0
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_lines: list[str] = []
        for row in sheet.iter_row():
            vals = [str(cell.value or "") for cell in row]
            sheet_lines.append(" | ".join(vals))
        if sheet_lines:
            text = "\n".join(sheet_lines)
            char_count += len(text)
            lines.append(f"\n--- Sheet: {sheet_name} ({len(sheet_lines)} rows) ---\n{text[:3000]}")
            if len(text) > 3000:
                lines.append(f"[... truncated, {len(text)} chars total in this sheet]")
    wb.close()
    lines.append(f"\n--- Extracted {char_count} chars from {len(wb.sheetnames)} sheets ---")
    return "\n".join(lines)


async def video_info(filepath: str) -> str:
    path = _confine(filepath)
    if not path.exists():
        return f"[error] File not found: {filepath}"
    try:
        import ffmpeg

        probe = await ffmpeg.probe(str(path))
    except ImportError:
        return "[error] ffmpeg-python not installed (pip install raven-agent[media])"
    except Exception as e:
        return f"[error] ffprobe failed: {e}"

    streams = probe.get("streams", [])
    info = probe.get("format", {})
    lines = [
        f"Video: {path.name}",
        f"- Duration: {info.get('duration', 'N/A')}s",
        f"- Size: {info.get('size', 'N/A')} bytes",
        f"- Bitrate: {info.get('bit_rate', 'N/A')} bps",
        f"- Format: {info.get('format_name', 'N/A')}",
        f"- Streams: {len(streams)}",
    ]
    for i, s in enumerate(streams):
        codec = s.get("codec_name", "?")
        stype = s.get("codec_type", "?")
        if stype == "video":
            lines.append(
                f"  Stream {i}: {stype} {codec}, {s.get('width', '?')}x{s.get('height', '?')}, {s.get('r_frame_rate', '?')} fps"
            )
        elif stype == "audio":
            lines.append(f"  Stream {i}: {stype} {codec}, {s.get('sample_rate', '?')} Hz, {s.get('channels', '?')} ch")
        else:
            lines.append(f"  Stream {i}: {stype} {codec}")
    return "\n".join(lines)


async def video_thumbnail(filepath: str, time_sec: float = 1.0, size: str = "320x240", output: str = "") -> str:
    path = _confine(filepath)
    exists = await asyncio.to_thread(path.exists)
    if not exists:
        return f"[error] File not found: {filepath}"
    try:
        import ffmpeg
    except ImportError:
        return "[error] ffmpeg-python not installed (pip install raven-agent[media])"

    if output:
        out_path = Path(output)
    else:
        out_path = path.parent / f"{path.stem}_thumb.jpg"
    await asyncio.to_thread(out_path.parent.mkdir, parents=True, exist_ok=True)
    try:
        await asyncio.to_thread(
            lambda: (
                ffmpeg.input(str(path), ss=time_sec)
                .filter("scale", *size.split("x"))
                .output(str(out_path), vframes=1)
                .overwrite_output()
                .run(capture_stdout=True, capture_stderr=True)
            )
        )
    except Exception as e:
        return f"[error] Thumbnail extraction failed: {e}"

    try:
        from PIL import Image

        def _encode() -> str:
            thumb = Image.open(out_path)
            buf = BytesIO()
            thumb.save(buf, format="JPEG", quality=85)
            return base64.b64encode(buf.getvalue()).decode()

        b64 = await asyncio.to_thread(_encode)
    except Exception as e:
        logger.debug("Thumbnail base64 encoding failed: {}", e)
        b64 = ""

    stat = await asyncio.to_thread(out_path.stat)
    lines = [
        f"Thumbnail saved: {out_path} ({stat.st_size} bytes)",
        f"- Time: {time_sec}s, Size: {size}",
    ]
    if b64:
        lines.append(f"![Thumbnail](data:image/jpeg;base64,{b64})")
    return "\n".join(lines)


async def video_transcribe(filepath: str, model: str = "whisper-1", language: str = "") -> str:
    """Extract audio from video and transcribe using OpenAI Whisper API."""
    path = _confine(filepath)
    if not path.exists():
        return f"[error] File not found: {filepath}"
    api_key = _OPENAI_API_KEY
    if not api_key:
        return "[error] OPENAI_API_KEY env var required for transcription"
    try:
        import httpx
    except ImportError:
        return "[error] httpx not available"

    tmp_wav = ""
    try:
        tmp_wav = str(Path(tempfile.gettempdir()) / f"raven_audio_{os.urandom(4).hex()}.wav")
        ffmpeg_cmd = [
            "ffmpeg",
            "-i",
            str(path),
            "-vn",
            "-acodec",
            "pcm_s16le",
            "-ar",
            "16000",
            "-ac",
            "1",
            "-y",
            tmp_wav,
        ]
        proc = await asyncio.create_subprocess_exec(
            *ffmpeg_cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        await proc.communicate()
        tmp_exists = await asyncio.to_thread(os.path.exists, tmp_wav)
        if not tmp_exists:
            return "[error] ffmpeg audio extraction failed (ffmpeg not found or invalid file)"

        def _read_audio() -> bytes:
            with Path(tmp_wav).open("rb") as f:
                return f.read()

        audio_data = await asyncio.to_thread(_read_audio)
        data = {"model": model}
        if language:
            data["language"] = language
        files = {"file": (path.stem + ".wav", audio_data, "audio/wav")}
        async with httpx.AsyncClient(timeout=300) as client:
            resp = await client.post(
                "https://api.openai.com/v1/audio/transcriptions",
                headers={"Authorization": f"Bearer {api_key}"},
                data=data,
                files=files,
            )
            resp.raise_for_status()
            result = resp.json()
            text = result.get("text", "")
            duration_sec = result.get("duration", 0)
            lines = [
                f"Transcription: {path.name}",
                f"Model: {model}",
                f"Duration: {duration_sec}s" if duration_sec else "",
                f"---\n{text}",
            ]
            return "\n".join(filter(None, lines))
    except Exception as e:
        logger.error("Video transcription failed: {}", e)
        return f"[error] Transcription failed: {e}"
    finally:
        if tmp_wav and Path(tmp_wav).exists():
            Path(tmp_wav).unlink()


async def video_extract_frames(
    filepath: str, interval_sec: float = 5.0, max_frames: int = 10, size: str = "640x480", output_dir: str = ""
) -> str:
    path = _confine(filepath)
    exists = await asyncio.to_thread(path.exists)
    if not exists:
        return f"[error] File not found: {filepath}"

    try:
        import ffmpeg
    except ImportError:
        return "[error] ffmpeg-python not installed (pip install raven-agent[media])"

    out_dir = Path(output_dir) if output_dir else (path.parent / f"{path.stem}_frames")
    await asyncio.to_thread(out_dir.mkdir, parents=True, exist_ok=True)
    try:
        probe = await asyncio.to_thread(lambda: ffmpeg.probe(str(path)))
        duration = float(probe["format"].get("duration", 0))
        if duration <= 0:
            return "[error] Could not determine video duration"
        n_frames = min(int(duration // interval_sec), max_frames)
        w, h = size.split("x") if "x" in size else ("640", "480")
        frame_paths = []
        for i in range(n_frames):
            ts = i * interval_sec
            out_path = str(out_dir / f"frame_{i:04d}_{int(ts)}s.jpg")

            def _extract_frame(p: Path = path, t: float = ts, o: str = out_path, w_in: str = w, h_in: str = h) -> None:
                ffmpeg.input(str(p), ss=t).filter("scale", w_in, h_in).output(o, vframes=1).overwrite_output().run(
                    capture_stdout=True, capture_stderr=True
                )

            await asyncio.to_thread(_extract_frame)
            frame_paths.append(out_path)
        previews = []
        for fp in frame_paths[:5]:
            try:

                def _read_frame(fp: str = fp) -> str:
                    with Path(fp).open("rb") as f:
                        return base64.b64encode(f.read()).decode()

                b64 = await asyncio.to_thread(_read_frame)
                previews.append(f"data:image/jpeg;base64,{b64}")
            except Exception:
                logger.warning("Failed to generate preview for frame")
                continue
        lines = [
            f"Extracted {len(frame_paths)} frames from {path.name}",
            f"Interval: {interval_sec}s, Size: {size}",
            f"Output: {out_dir}",
        ]
        for p in previews:
            lines.append(f"![frame]({p})")
        return "\n".join(lines)
    except Exception as e:
        logger.error("Frame extraction failed: {}", e)
        return f"[error] Frame extraction failed: {e}"


async def image_analyze(filepath: str, prompt: str = "Describe this image in detail") -> str:
    path = _confine(filepath)
    exists = await asyncio.to_thread(path.exists)
    if not exists:
        return f"[error] File not found: {filepath}"
    api_key = _OPENAI_API_KEY
    if not api_key:
        return "[error] OPENAI_API_KEY env var required for image analysis"
    try:
        import httpx

        def _read_image() -> str:
            with Path(path).open("rb") as f:
                return base64.b64encode(f.read()).decode()

        b64 = await asyncio.to_thread(_read_image)
        ext = path.suffix.lower().lstrip(".")
        if ext in ("jpg", "jpeg"):
            mime = "image/jpeg"
        elif ext == "png":
            mime = "image/png"
        elif ext == "gif":
            mime = "image/gif"
        elif ext == "webp":
            mime = "image/webp"
        else:
            return "[error] Unsupported image format. Supported: jpg, png, gif, webp"
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                "https://api.openai.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={
                    "model": "gpt-4o-mini",
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": prompt},
                                {
                                    "type": "image_url",
                                    "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "low"},
                                },
                            ],
                        },
                    ],
                    "max_tokens": 1000,
                },
            )
            data: dict[str, Any] = r.json()
            if "choices" in data and len(data["choices"]) > 0:
                return str(data["choices"][0]["message"]["content"])
            return f"[error] API response: {data}"
    except Exception as e:
        logger.error("Image analysis failed: {}", e)
        return f"[error] Image analysis failed: {e}"


def register_media_tools(registry: ToolRegistry) -> None:
    registry.register(
        ToolSpec(
            name="image_generate",
            description="Generate an image from a text prompt using DALL-E 3 or Stable Diffusion",
            parameters={
                "prompt": {"type": "string", "description": "Image description", "required": True},
                "size": {
                    "type": "string",
                    "description": "Image size: 1024x1024, 1792x1024, or 1024x1792",
                    "required": False,
                },
                "quality": {"type": "string", "description": "standard or hd", "required": False},
                "n": {"type": "integer", "description": "Number of images (1)", "required": False},
                "model": {
                    "type": "string",
                    "description": "Model: dall-e-3 (default) or stable-diffusion",
                    "required": False,
                },
            },
            handler=image_generate,
            category="media",
            timeout=120,
        )
    )
    registry.register(
        ToolSpec(
            name="image_edit",
            description="Edit an image file: resize, crop, rotate, flip, convert format",
            parameters={
                "filepath": {"type": "string", "description": "Path to image file", "required": True},
                "resize": {"type": "string", "description": "New size as WxH (e.g. 800x600)", "required": False},
                "crop": {"type": "string", "description": "Crop rect as left,upper,right,lower", "required": False},
                "rotate": {"type": "integer", "description": "Rotation angle in degrees", "required": False},
                "flip": {"type": "string", "description": "horizontal or vertical", "required": False},
                "format": {"type": "string", "description": "Output format: JPEG, PNG, GIF, WEBP", "required": False},
                "quality": {"type": "integer", "description": "Output quality 1-100", "required": False},
                "output": {"type": "string", "description": "Output file path (optional)", "required": False},
            },
            handler=image_edit,
            category="media",
            timeout=60,
        )
    )
    registry.register(
        ToolSpec(
            name="image_analyze",
            description="Analyze an image using GPT-4o Vision (describe, caption, or answer questions about an image)",
            parameters={
                "filepath": {"type": "string", "description": "Path to image file", "required": True},
                "prompt": {"type": "string", "description": "Question or prompt about the image", "required": False},
            },
            handler=image_analyze,
            category="media",
            timeout=60,
        )
    )
    registry.register(
        ToolSpec(
            name="document_parse",
            description="Extract text content from documents (PDF, DOCX, PPTX, XLSX)",
            parameters={
                "filepath": {"type": "string", "description": "Path to document file", "required": True},
                "pages": {
                    "type": "string",
                    "description": "Page range for PDFs, e.g. 1-3,5 (optional)",
                    "required": False,
                },
            },
            handler=document_parse,
            category="media",
            timeout=60,
        )
    )
    registry.register(
        ToolSpec(
            name="video_transcribe",
            description="Transcribe audio from a video file using OpenAI Whisper API",
            parameters={
                "filepath": {"type": "string", "description": "Path to video file", "required": True},
                "model": {"type": "string", "description": "Whisper model (whisper-1)", "required": False},
                "language": {"type": "string", "description": "Language code (optional, e.g. en)", "required": False},
            },
            handler=video_transcribe,
            category="media",
            timeout=300,
        )
    )
    registry.register(
        ToolSpec(
            name="audio_transcribe",
            description="Transcribe audio from an audio file using OpenAI Whisper API",
            parameters={
                "filepath": {
                    "type": "string",
                    "description": "Path to audio file (mp3, wav, ogg, flac, m4a)",
                    "required": True,
                },
                "model": {"type": "string", "description": "Whisper model (whisper-1)", "required": False},
                "language": {"type": "string", "description": "Language code (optional, e.g. en)", "required": False},
            },
            handler=video_transcribe,
            category="media",
            timeout=300,
        )
    )
    registry.register(
        ToolSpec(
            name="video_extract_frames",
            description="Extract frames from a video at regular intervals as JPEG images",
            parameters={
                "filepath": {"type": "string", "description": "Path to video file", "required": True},
                "interval_sec": {
                    "type": "number",
                    "description": "Interval between frames in seconds (default 5.0)",
                    "required": False,
                },
                "max_frames": {
                    "type": "integer",
                    "description": "Maximum number of frames to extract (default 10)",
                    "required": False,
                },
                "size": {"type": "string", "description": "Frame size as WxH (default 640x480)", "required": False},
                "output_dir": {"type": "string", "description": "Output directory (optional)", "required": False},
            },
            handler=video_extract_frames,
            category="media",
            timeout=120,
        )
    )
    registry.register(
        ToolSpec(
            name="video_info",
            description="Get metadata information about a video file (duration, codec, resolution, streams)",
            parameters={
                "filepath": {"type": "string", "description": "Path to video file", "required": True},
            },
            handler=video_info,
            category="media",
            timeout=30,
        )
    )
    registry.register(
        ToolSpec(
            name="video_thumbnail",
            description="Extract a thumbnail image from a video at a specific timestamp",
            parameters={
                "filepath": {"type": "string", "description": "Path to video file", "required": True},
                "time_sec": {"type": "number", "description": "Timestamp in seconds", "required": False},
                "size": {"type": "string", "description": "Thumbnail size as WxH", "required": False},
                "output": {"type": "string", "description": "Output file path (optional)", "required": False},
            },
            handler=video_thumbnail,
            category="media",
            timeout=60,
        )
    )
